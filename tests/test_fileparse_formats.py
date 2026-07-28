from augur.knowledge import fileparse


def test_pdf_empty_password_owner_only_not_encrypted(tmp_path):
    from pypdf import PdfWriter

    path = tmp_path / "owner_only.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.encrypt(user_password="", owner_password="owner-only")
    with path.open("wb") as f:
        writer.write(f)

    text, reason = fileparse.extract_text(str(path))
    assert text is None
    assert reason in ("no_text", "pdf")
    assert reason != "encrypted"


def test_pdf_real_password_reports_encrypted(tmp_path):
    from pypdf import PdfWriter

    path = tmp_path / "locked.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.encrypt(user_password="need-pwd", owner_password="owner")
    with path.open("wb") as f:
        writer.write(f)

    text, reason = fileparse.extract_text(str(path))
    assert text is None
    assert reason == "encrypted"


def test_docx_extracts_paragraph_and_table(tmp_path):
    import docx

    path = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.add_paragraph("alpha")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "beta"
    table.rows[0].cells[1].text = "gamma"
    doc.save(path)

    text, reason = fileparse.extract_text(str(path))
    assert reason == "docx"
    assert "alpha" in text
    assert "beta\tgamma" in text


def test_xlsx_extracts_cells(tmp_path):
    import openpyxl

    path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "hello"
    ws["B1"] = 42
    wb.save(path)
    wb.close()

    text, reason = fileparse.extract_text(str(path))
    assert reason == "xlsx"
    assert "hello\t42" in text


def test_pptx_extracts_slide_text(tmp_path):
    from pptx import Presentation

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(left=0, top=0, width=prs.slide_width // 2, height=prs.slide_height // 4)
    box.text_frame.text = "deck text"
    prs.save(path)

    text, reason = fileparse.extract_text(str(path))
    assert reason == "pptx"
    assert "deck text" in text


def test_image_without_ocr_dependency_returns_missing_ocr(tmp_path, monkeypatch):
    path = tmp_path / "img.png"
    path.write_bytes(b"not-an-image")

    real_import = __import__

    def fake_import(name, *args, **kwargs):
        if name in ("pytesseract", "PIL", "PIL.Image", "PIL.ImageOps"):
            raise ImportError(name)
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    text, reason = fileparse.extract_text(str(path))
    assert text is None
    assert reason == "missing_ocr"


def test_doc_without_parser_returns_missing_parser(tmp_path, monkeypatch):
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"legacy-doc")

    monkeypatch.setattr(fileparse, "_soffice_bin", lambda: None)
    text, reason = fileparse.extract_text(str(path))
    assert text is None
    assert reason == "missing_parser"


def test_ppt_without_parser_returns_missing_parser(tmp_path, monkeypatch):
    path = tmp_path / "legacy.ppt"
    path.write_bytes(b"legacy-ppt")

    monkeypatch.setattr(fileparse, "_soffice_bin", lambda: None)
    text, reason = fileparse.extract_text(str(path))
    assert text is None
    assert reason == "missing_parser"


def test_xls_without_parser_returns_missing_parser(tmp_path, monkeypatch):
    path = tmp_path / "legacy.xls"
    path.write_bytes(b"legacy-xls")

    real_import = __import__

    def fake_import(name, *args, **kwargs):
        if name == "xlrd":
            raise ImportError(name)
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)
    monkeypatch.setattr(fileparse, "_soffice_bin", lambda: None)
    text, reason = fileparse.extract_text(str(path))
    assert text is None
    assert reason == "missing_parser"
