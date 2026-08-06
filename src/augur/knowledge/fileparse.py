"""本機檔多格式文字抽取器(admin 控制台 `+資料夾` 核心,計畫 §三)— 純規則、零 LLM/零 API。

🎯 這支在做什麼(白話):把本機任意副檔名檔案抽出**逐字文字**(不摘要、不改寫、不杜撰),
   供知識層切句/嵌入/檢索。每格式一個抽取器;抽不出(未知/二進位/損壞/加密/超大)= **誠實跳過並分類記數**,
   絕不硬湊內容(#1 逐字 · #15 誠實)。惡意檔防護:大小上限(防 OOM)、符號連結不跟(防逃逸)、
   解析例外吞為 skip(防單檔崩整批)。
   PDF-C：`extract_text(..., ocr_pdf=True)` 對弱／空字層 PDF 走 Tesseract（預設關；禁 ASR／caption）。
守 #1(逐字、禁 AI 生成/改寫)· #15(抽不出誠實跳過、不杜撰)· #5(惡意檔/大小/符號連結防護)· #29。

執行指令矩陣(本檔=library;CLI 見 scripts/acquire_local_files.py / backfill_pdf_ocr.py):
  python -c "from augur.knowledge.fileparse import extract_text; print(extract_text('X.pdf'))"
  python -c "from augur.knowledge.fileparse import extract_text; print(extract_text('X.pdf', ocr_pdf=True))"

自測（本檔=library #18；免 DB 免 API 可個別驗證）：
  python -m augur.knowledge.fileparse              # 印用途+公開入口（唯讀）
  python -m augur.knowledge.fileparse --selftest   # 結構＋暫存 PDF 加密假陽性／真密碼路徑
"""
from __future__ import annotations

import os
import shutil
import subprocess
import tempfile

MAX_BYTES = 50 * 1024 * 1024      # 單檔上限 50MB(防 OOM/zip bomb 解壓;超過=oversize skip)
_TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json", ".xml",
             ".yaml", ".yml", ".rst", ".ini", ".toml",
             ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".go", ".rs", ".sh", ".sql", ".r"}
_IMAGE_EXT = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".tif", ".tiff", ".bmp"}
_OLE_OFFICE_EXT = {".doc", ".xls", ".ppt"}
# skip 分類(誠實記數,計畫 §三):
SKIP = ("oversize", "symlink", "empty", "decode_error", "parse_error",
        "encrypted", "unknown_ext", "no_text", "missing_ocr", "missing_parser",
        "ocr_quality")
# 人類可讀（UI／摘要；reason key 仍用 SKIP 英文碼）
SKIP_LABEL_ZH = {
    "oversize": "超過單檔大小上限",
    "symlink": "符號連結已略過",
    "empty": "空檔",
    "decode_error": "文字解碼失敗",
    "parse_error": "解析失敗／損壞",
    "encrypted": "加密 PDF，需密碼／請先解密",
    "unknown_ext": "未支援副檔名",
    "no_text": "無可抽文字（如掃描圖）",
    "missing_ocr": "缺 OCR 引擎（需 pytesseract+tesseract）",
    "missing_parser": "缺解析器或系統轉檔器",
    "ocr_quality": "OCR 品質未過閘（過短／空白頁過多）",
}

# PDF-C OCR（code-go）：預設關；extract_text(..., ocr_pdf=True) 或 backfill 才開
OCR_TRIGGER_CHARS_DEFAULT = 200
OCR_MAX_PAGES_DEFAULT = 40
OCR_MIN_CHARS = 80          # 弱於 pypdf 字層則寧可保留字層／跳過
OCR_BLANK_PAGE_RATIO_MAX = 0.50
OCR_DPI = 200
OCR_LANGS = "chi_tra+eng"
S0_OCR_MARK = "<!-- via=pdf_ocr -->\n"   # source_mark=S0；仍用 local_upload


def _read_text(path):
    """純文字檔:chardet 偵編碼、errors='replace' 兜底(仿 fetch_oa_fulltext 慣例)。"""
    raw = open(path, "rb").read()
    if not raw:
        return None, "empty"
    enc = "utf-8"
    try:
        import chardet
        d = chardet.detect(raw[:65536])
        enc = d.get("encoding") or "utf-8"
    except Exception:
        pass
    try:
        return raw.decode(enc, errors="replace"), "text"
    except Exception:
        return raw.decode("utf-8", errors="replace"), "text"


def _read_pdf(path):
    """PDF 逐字抽取。is_encrypted 常為「僅 owner 限制」假陽性——先試空密碼再開；真需密碼才 skip encrypted。"""
    from pypdf import PdfReader
    r = PdfReader(path)
    if getattr(r, "is_encrypted", False):
        try:
            # pypdf: 0=NOT_DECRYPTED（真需密碼）; 1/2=USER/OWNER（含空 user 密碼之 owner 限制檔）
            if not r.decrypt(""):
                return None, "encrypted"
        except Exception:
            return None, "encrypted"
    txt = "\n".join((p.extract_text() or "") for p in r.pages)
    return (txt, "pdf") if txt.strip() else (None, "no_text")   # 弱／空字層可走 ocr_pdf


def ocr_pdf_pages(
    path,
    *,
    max_pages: int = OCR_MAX_PAGES_DEFAULT,
    dpi: int = OCR_DPI,
    langs: str = OCR_LANGS,
):
    """PDF 頁光柵 → Tesseract。回 (text|None, reason, meta:dict)。

    reason 成功＝`pdf_ocr`；失敗∈ SKIP。禁 ASR／caption；不經 LLM（PDF-C）。
    """
    meta = {
        "n_pages": 0, "n_ocr_pages": 0, "blank_pages": 0,
        "chars": 0, "truncated": False, "dpi": dpi, "langs": langs,
    }
    try:
        import fitz  # pymupdf
        import pytesseract
        from PIL import Image
    except Exception:
        return None, "missing_ocr", meta
    try:
        doc = fitz.open(path)
    except Exception:
        return None, "parse_error", meta
    try:
        n_all = len(doc)
        meta["n_pages"] = n_all
        limit = max(1, int(max_pages))
        if n_all > limit:
            meta["truncated"] = True
        n = min(n_all, limit)
        parts = []
        blank = 0
        for i in range(n):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=dpi)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            try:
                s = pytesseract.image_to_string(img, lang=langs) or ""
            except pytesseract.TesseractNotFoundError:
                return None, "missing_ocr", meta
            except Exception:
                return None, "parse_error", meta
            if not s.strip():
                blank += 1
            parts.append(s)
            meta["n_ocr_pages"] = i + 1
        meta["blank_pages"] = blank
        if n > 0 and (blank / n) > OCR_BLANK_PAGE_RATIO_MAX:
            return None, "ocr_quality", meta
        text = "\n".join(parts).strip()
        meta["chars"] = len(text)
        if len(text) < OCR_MIN_CHARS:
            return None, "ocr_quality", meta
        return text, "pdf_ocr", meta
    finally:
        try:
            doc.close()
        except Exception:
            pass


def _pdf_needs_ocr(text, reason, *, trigger_chars: int) -> bool:
    if reason == "encrypted":
        return False
    if reason == "no_text" or text is None:
        return True
    return len((text or "").strip()) < int(trigger_chars)


def _read_docx(path):
    import docx
    d = docx.Document(path)
    lines = [p.text for p in d.paragraphs]
    for tbl in d.tables:                                   # 表格逐列(docstring 承諾「段落+表格」;paragraphs 不含表格)
        for row in tbl.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                lines.append("\t".join(cells))
    txt = "\n".join(lines)
    return (txt, "docx") if txt.strip() else (None, "no_text")


def _read_pptx(path):
    from pptx import Presentation
    p = Presentation(path)
    txt = "\n".join(sh.text for sl in p.slides for sh in sl.shapes if sh.has_text_frame)
    return (txt, "pptx") if txt.strip() else (None, "no_text")


def _read_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    wb.close()
    txt = "\n".join(lines)
    return (txt, "xlsx") if txt.strip() else (None, "no_text")


def _soffice_bin():
    return shutil.which("soffice") or shutil.which("libreoffice")


def _read_via_soffice(path, out_ext, method):
    soffice = _soffice_bin()
    if not soffice:
        return None, "missing_parser"
    with tempfile.TemporaryDirectory() as td:
        cmd = [soffice, "--headless", "--convert-to", out_ext.lstrip("."), "--outdir", td, path]
        try:
            subprocess.run(cmd, capture_output=True, text=True, timeout=60, check=True)
        except (OSError, subprocess.SubprocessError):
            return None, "parse_error"
        out = os.path.join(td, os.path.splitext(os.path.basename(path))[0] + out_ext)
        if not os.path.isfile(out):
            return None, "parse_error"
        if out_ext == ".txt":
            txt, status = _read_text(out)
            if txt is None:
                return None, status
            return (txt, method) if txt.strip() else (None, "no_text")
        return extract_text(out)


def _read_doc(path):
    return _read_via_soffice(path, ".txt", "doc")


def _read_ppt(path):
    return _read_via_soffice(path, ".txt", "ppt")


def _read_xls(path):
    try:
        import xlrd
    except Exception:
        return _read_via_soffice(path, ".txt", "xls")
    book = xlrd.open_workbook(path, on_demand=True)
    lines = []
    try:
        for sheet in book.sheets():
            for rx in range(sheet.nrows):
                cells = [str(v) for v in sheet.row_values(rx) if str(v).strip()]
                if cells:
                    lines.append("\t".join(cells))
    finally:
        book.release_resources()
    txt = "\n".join(lines)
    return (txt, "xls") if txt.strip() else (None, "no_text")


def _read_image(path):
    try:
        import pytesseract
        from PIL import Image, ImageOps
    except Exception:
        return None, "missing_ocr"
    try:
        with Image.open(path) as im:
            im = ImageOps.exif_transpose(im)
            if im.mode not in ("L", "RGB"):
                im = im.convert("RGB")
            text = pytesseract.image_to_string(im)
    except pytesseract.TesseractNotFoundError:
        return None, "missing_ocr"
    except Exception:
        return None, "parse_error"
    return (text, "image_ocr") if text.strip() else (None, "no_text")


def _read_epub(path):
    import re
    import html as _html
    from ebooklib import epub, ITEM_DOCUMENT
    b = epub.read_epub(path)
    parts = []
    for it in b.get_items_of_type(ITEM_DOCUMENT):
        raw = it.get_content().decode("utf-8", "replace")
        raw = re.sub(r"(?is)<(script|style).*?</\1>", " ", raw)
        parts.append(_html.unescape(re.sub(r"(?s)<[^>]+>", " ", raw)))
    txt = re.sub(r"\s+\n", "\n", "\n".join(parts))
    return (txt, "epub") if txt.strip() else (None, "no_text")


def extract_text(
    path,
    *,
    ocr_pdf: bool = False,
    ocr_trigger_chars: int | None = None,
    ocr_max_pages: int | None = None,
):
    """回 (text|None, method_or_skipreason)。text 非 None = 抽取成功;None = 跳過(reason ∈ SKIP)。

    純規則、逐字;任何解析例外 → ('parse_error') 不外拋(單檔不崩整批,#5/#15)。
    ocr_pdf：PDF 字層弱／空時用 Tesseract 補抽（預設 False；PDF-C；禁 ASR／caption）。"""
    try:
        if os.path.islink(path):
            return None, "symlink"                 # 不跟符號連結(防逃逸/迴圈)
        size = os.path.getsize(path)
        if size == 0:
            return None, "empty"
        if size > MAX_BYTES:
            return None, "oversize"
        ext = os.path.splitext(path)[1].lower()
        if ext in _TEXT_EXT:
            return _read_text(path)
        if ext == ".pdf":
            text, reason = _read_pdf(path)
            trigger = (ocr_trigger_chars if ocr_trigger_chars is not None
                       else OCR_TRIGGER_CHARS_DEFAULT)
            max_pages = (ocr_max_pages if ocr_max_pages is not None
                         else OCR_MAX_PAGES_DEFAULT)
            if ocr_pdf and _pdf_needs_ocr(text, reason, trigger_chars=trigger):
                ocr_text, ocr_reason, _meta = ocr_pdf_pages(path, max_pages=max_pages)
                if ocr_reason == "pdf_ocr" and ocr_text:
                    if len(ocr_text.strip()) > len((text or "").strip()):
                        return ocr_text, "pdf_ocr"
                elif text is None and ocr_reason in SKIP:
                    return None, ocr_reason
            return (text, reason) if text is not None else (None, reason)
        if ext == ".docx":
            return _read_docx(path)
        if ext == ".doc":
            return _read_doc(path)
        if ext == ".pptx":
            return _read_pptx(path)
        if ext == ".ppt":
            return _read_ppt(path)
        if ext == ".xlsx":
            return _read_xlsx(path)
        if ext == ".xls":
            return _read_xls(path)
        if ext in _IMAGE_EXT:
            return _read_image(path)
        if ext == ".epub":
            return _read_epub(path)
        if ext in (".html", ".htm"):
            raw, st = _read_text(path)
            if not raw:
                return None, st
            try:
                import re
                import html as _html
                # stdlib 剝標(與 fetch_oa_fulltext.strip_html 同法、零依賴)
                t = re.sub(r"(?is)<(script|style).*?</\1>", " ", raw)
                t = re.sub(r"(?s)<[^>]+>", " ", t)
                t = _html.unescape(t)
                t = re.sub(r"\s+\n", "\n", t)
                return (t, "html") if t.strip() else (None, "no_text")
            except Exception:
                return None, "parse_error"
        return None, "unknown_ext"                 # 二進位/未支援 → 誠實跳過(不杜撰)
    except Exception:
        return None, "parse_error"                 # 損壞/權限/解析崩 → skip(不崩整批)


def _selftest():
    """結構不變式＋暫存 PDF：owner 限制（空密碼）可抽；真 user 密碼 → encrypted。"""
    import tempfile
    ok = True
    def chk(name, cond):
        nonlocal ok; ok = ok and cond
        print(f"  {'✓' if cond else '✗FAIL'} {name}")
    chk("MAX_BYTES=50MB", MAX_BYTES == 50 * 1024 * 1024)
    chk("SKIP 全十一類齊備", set(SKIP) == {"oversize", "symlink", "empty", "decode_error",
                                        "parse_error", "encrypted", "unknown_ext", "no_text",
                                        "missing_ocr", "missing_parser", "ocr_quality"})
    chk("SKIP_LABEL_ZH 覆蓋 SKIP", set(SKIP_LABEL_ZH) == set(SKIP))
    chk("encrypted 標籤含密碼提示", "密碼" in SKIP_LABEL_ZH["encrypted"])
    chk(".pdf 不在 _TEXT_EXT(走專屬抽取器)", ".pdf" not in _TEXT_EXT)
    chk(".txt/.md/.csv/.json 屬純文字集", {".txt", ".md", ".csv", ".json"} <= _TEXT_EXT)
    chk("常見 image 副檔名納入 OCR 集", {".jpg", ".png", ".webp", ".gif", ".tiff", ".bmp"} <= _IMAGE_EXT)
    chk("舊版 Office 副檔名納入辨識", {".doc", ".xls", ".ppt"} <= _OLE_OFFICE_EXT)
    chk("公開入口 extract_text 存在", callable(extract_text))
    chk("ocr_pdf_pages 公開", callable(ocr_pdf_pages))
    chk("needs_ocr 弱字層", _pdf_needs_ocr("x" * 50, "pdf", trigger_chars=200) is True)
    chk("needs_ocr 足字層否", _pdf_needs_ocr("x" * 500, "pdf", trigger_chars=200) is False)
    chk("needs_ocr encrypted 否", _pdf_needs_ocr(None, "encrypted", trigger_chars=200) is False)
    chk("S0 mark", S0_OCR_MARK.startswith("<!-- via=pdf_ocr"))
    chk("各格式抽取器齊備", all(callable(f) for f in
                              (_read_text, _read_pdf, _read_docx, _read_doc, _read_pptx,
                               _read_ppt, _read_xlsx, _read_xls, _read_image, _read_epub)))
    with tempfile.TemporaryDirectory() as td:
        png = os.path.join(td, "broken.png")
        with open(png, "wb") as f:
            f.write(b"not-an-image")
        _text, why = extract_text(png)
        chk("image 不誤落 unknown_ext", why in ("missing_ocr", "parse_error"))
    try:
        from pypdf import PdfWriter
        with tempfile.TemporaryDirectory() as td:
            owner_path = os.path.join(td, "owner_restrict.pdf")
            w = PdfWriter()
            w.add_blank_page(width=200, height=200)
            w.encrypt(user_password="", owner_password="owner-only")
            with open(owner_path, "wb") as f:
                w.write(f)
            text, reason = _read_pdf(owner_path)
            # 空白頁無文字層 → no_text（非 encrypted）；關鍵＝不誤判 encrypted
            chk("owner 限制空密碼不標 encrypted", reason != "encrypted")
            chk("owner 限制回 no_text 或 pdf", reason in ("no_text", "pdf"))

            locked_path = os.path.join(td, "user_locked.pdf")
            w2 = PdfWriter()
            w2.add_blank_page(width=200, height=200)
            w2.encrypt(user_password="need-pwd", owner_password="owner")
            with open(locked_path, "wb") as f:
                w2.write(f)
            text2, reason2 = _read_pdf(locked_path)
            chk("真 user 密碼 → encrypted", text2 is None and reason2 == "encrypted")
            text3, reason3 = extract_text(locked_path)
            chk("extract_text 真加密仍 skip encrypted", text3 is None and reason3 == "encrypted")
    except ImportError:
        print("  · SKIP PDF 加解密測（無 pypdf）")
    except Exception as e:
        chk(f"PDF 加解密測無例外({type(e).__name__})", False)
        print(f"    detail: {e}")
    print("自測:" + ("全通過 ✓" if ok else "有 FAIL ✗"))
    return 0 if ok else 1


if __name__ == "__main__":
    import sys
    if "--selftest" in sys.argv:
        sys.exit(_selftest())
    print((__doc__ or __name__).split("🎯")[0].strip())
    print("(自測:python -m augur.knowledge.fileparse --selftest;免 DB 免 API)")
